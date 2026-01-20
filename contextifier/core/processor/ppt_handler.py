# libs/core/processor/ppt_handler.py
"""
PPT Handler - PPT/PPTX Document Processor

Class-based handler for PPT/PPTX files inheriting from BaseHandler.
"""
import logging
from typing import Any, Dict, List, Optional, Set, TYPE_CHECKING

from pptx import Presentation

from contextifier.core.processor.base_handler import BaseHandler
from contextifier.core.processor.ppt_helper import (
    ElementType,
    SlideElement,
    extract_ppt_metadata,
    format_metadata,
    extract_text_with_bullets,
    is_simple_table,
    extract_simple_table_as_text,
    convert_table_to_html,
    extract_table_as_text,
    extract_chart_data,
    get_shape_position,
    is_picture_shape,
    process_image_shape,
    process_group_shape,
    extract_slide_notes,
    merge_slide_elements,
)

if TYPE_CHECKING:
    from contextifier.core.document_processor import CurrentFile

logger = logging.getLogger("document-processor")


class PPTHandler(BaseHandler):
    """PPT/PPTX File Processing Handler Class"""
    
    def extract_text(
        self,
        current_file: "CurrentFile",
        extract_metadata: bool = True,
        **kwargs
    ) -> str:
        """
        Extract text from PPT/PPTX file.
        
        Args:
            current_file: CurrentFile dict containing file info and binary data
            extract_metadata: Whether to extract metadata
            **kwargs: Additional options
            
        Returns:
            Extracted text
        """
        file_path = current_file.get("file_path", "unknown")
        self.logger.info(f"PPT processing: {file_path}")
        return self._extract_ppt_enhanced(current_file, extract_metadata)
    
    def _extract_ppt_enhanced(self, current_file: "CurrentFile", extract_metadata: bool = True) -> str:
        """Enhanced PPT processing."""
        file_path = current_file.get("file_path", "unknown")
        self.logger.info(f"Enhanced PPT processing: {file_path}")
        
        try:
            # Open from stream to avoid path encoding issues
            file_stream = self.get_file_stream(current_file)
            prs = Presentation(file_stream)
            result_parts = []
            processed_images: Set[str] = set()
            total_tables = 0
            total_images = 0
            
            if extract_metadata:
                metadata = extract_ppt_metadata(prs)
                metadata_text = format_metadata(metadata)
                if metadata_text:
                    result_parts.append(metadata_text)
                    result_parts.append("")
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_tag = self.create_slide_tag(slide_idx + 1)
                result_parts.append(f"\n{slide_tag}\n")
                
                elements: List[SlideElement] = []
                
                for shape in slide.shapes:
                    try:
                        position = get_shape_position(shape)
                        shape_id = shape.shape_id if hasattr(shape, 'shape_id') else id(shape)
                        
                        if shape.has_table:
                            if is_simple_table(shape.table):
                                simple_text = extract_simple_table_as_text(shape.table)
                                if simple_text:
                                    elements.append(SlideElement(
                                        element_type=ElementType.TEXT,
                                        content=simple_text,
                                        position=position,
                                        shape_id=shape_id
                                    ))
                            else:
                                table_html = convert_table_to_html(shape.table)
                                if table_html:
                                    total_tables += 1
                                    elements.append(SlideElement(
                                        element_type=ElementType.TABLE,
                                        content=table_html,
                                        position=position,
                                        shape_id=shape_id
                                    ))
                        
                        elif is_picture_shape(shape):
                            image_tag = process_image_shape(shape, processed_images, self.image_processor)
                            if image_tag:
                                total_images += 1
                                elements.append(SlideElement(
                                    element_type=ElementType.IMAGE,
                                    content=image_tag,
                                    position=position,
                                    shape_id=shape_id
                                ))
                        
                        elif shape.has_chart:
                            chart_text = extract_chart_data(shape.chart)
                            if chart_text:
                                elements.append(SlideElement(
                                    element_type=ElementType.CHART,
                                    content=chart_text,
                                    position=position,
                                    shape_id=shape_id
                                ))
                        
                        elif hasattr(shape, "text_frame") and shape.text_frame:
                            text_content = extract_text_with_bullets(shape.text_frame)
                            if text_content:
                                elements.append(SlideElement(
                                    element_type=ElementType.TEXT,
                                    content=text_content,
                                    position=position,
                                    shape_id=shape_id
                                ))
                        
                        elif hasattr(shape, "text") and shape.text.strip():
                            elements.append(SlideElement(
                                element_type=ElementType.TEXT,
                                content=shape.text.strip(),
                                position=position,
                                shape_id=shape_id
                            ))
                        
                        elif hasattr(shape, "shapes"):
                            group_elements = process_group_shape(shape, processed_images, self.image_processor)
                            elements.extend(group_elements)
                    
                    except Exception as shape_e:
                        self.logger.warning(f"Error processing shape in slide {slide_idx + 1}: {shape_e}")
                        continue
                
                elements.sort(key=lambda e: e.sort_key)
                slide_content = merge_slide_elements(elements)
                
                if slide_content.strip():
                    result_parts.append(slide_content)
                else:
                    result_parts.append("[Empty Slide]\n")
                
                notes_text = extract_slide_notes(slide)
                if notes_text:
                    result_parts.append(f"\n[Slide Notes]\n{notes_text}\n")
            
            result = "".join(result_parts)
            self.logger.info(f"Enhanced PPT: {len(prs.slides)} slides, {total_tables} tables, {total_images} images")
            
            return result
            
        except Exception as e:
            self.logger.error(f"Error in enhanced PPT processing: {e}")
            import traceback
            self.logger.debug(traceback.format_exc())
            return self._extract_ppt_simple(current_file)
    
    def _extract_ppt_simple(self, current_file: "CurrentFile") -> str:
        """Simple text extraction (fallback)."""
        try:
            file_stream = self.get_file_stream(current_file)
            prs = Presentation(file_stream)
            result_parts = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_tag = self.create_slide_tag(slide_idx + 1)
                result_parts.append(f"\n{slide_tag}\n")
                
                slide_texts = []
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                        elif hasattr(shape, "table"):
                            table_text = extract_table_as_text(shape.table)
                            if table_text:
                                slide_texts.append(table_text)
                    except:
                        continue
                
                if slide_texts:
                    result_parts.append("\n".join(slide_texts) + "\n")
                else:
                    result_parts.append("[Empty Slide]\n")
            
            return "".join(result_parts)
            
        except Exception as e:
            self.logger.error(f"Error in simple PPT extraction: {e}")
            return f"[PPT file processing failed: {str(e)}]"
