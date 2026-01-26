# contextifier/core/processor/ppt_helper/ppt_table_extractor.py
"""
PPT Table Extractor - PPT/PPTX Format-Specific Table Extraction

Implements table extraction for PPT/PPTX files (python-pptx).
Follows BaseTableExtractor interface from table_extractor.py.

PPTX Table Structure:
- python-pptx Table object with rows and columns
- Cell merge via is_merge_origin, span_width, span_height
- Or via XML attributes: gridSpan, rowSpan

2-Pass Approach:
1. Pass 1: Detect table regions (shapes with has_table attribute)
2. Pass 2: Extract content from detected regions (TableData objects)

Usage:
    from contextifier.core.processor.ppt_helper.ppt_table_extractor import (
        PPTTableExtractor,
    )

    extractor = PPTTableExtractor()
    tables = extractor.extract_tables(presentation)  # python-pptx Presentation
"""
import logging
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple

from contextifier.core.functions.table_extractor import (
    BaseTableExtractor,
    TableCell,
    TableData,
    TableRegion,
    TableExtractorConfig,
)

logger = logging.getLogger("document-processor")


# ============================================================================
# Configuration
# ============================================================================

@dataclass
class PPTTableExtractorConfig(TableExtractorConfig):
    """Configuration specific to PPT table extraction.
    
    Attributes:
        skip_single_cell_tables: Whether to skip 1x1 tables (often containers)
        skip_single_column_tables: Whether to skip single column (Nx1) tables
        skip_single_row_tables: Whether to skip single row (1xN) tables
        treat_first_row_as_header: Whether to treat first row as header
    """
    skip_single_cell_tables: bool = True
    skip_single_column_tables: bool = True
    skip_single_row_tables: bool = True
    treat_first_row_as_header: bool = True


# ============================================================================
# Data Classes for Table Region Tracking
# ============================================================================

@dataclass
class PPTTableRegionInfo:
    """Additional information for PPT table region.
    
    Stores reference to the actual table object for Pass 2.
    """
    slide_index: int  # 0-based slide index
    shape_id: int     # Shape ID containing the table
    table_ref: Any    # Reference to python-pptx Table object


# ============================================================================
# PPTTableExtractor Class
# ============================================================================

class PPTTableExtractor(BaseTableExtractor):
    """PPT/PPTX format-specific table extractor.
    
    Extracts tables from PPT/PPTX files using python-pptx.
    Implements BaseTableExtractor interface.
    
    Supports:
    - Cell merges (colspan via gridSpan/span_width, rowspan via rowSpan/span_height)
    - Header row detection
    - Simple table detection (1xN, Nx1)
    """
    
    def __init__(self, config: Optional[PPTTableExtractorConfig] = None):
        """Initialize PPT table extractor.
        
        Args:
            config: PPT table extraction configuration
        """
        self._config = config or PPTTableExtractorConfig()
        super().__init__(self._config)
        # Cache for presentation and tables
        self._presentation_cache: Any = None
        self._table_info_cache: Dict[int, PPTTableRegionInfo] = {}
    
    def supports_format(self, format_type: str) -> bool:
        """Check if this extractor supports the given format."""
        return format_type.lower() in ("ppt", "pptx", "pptm")
    
    # ========================================================================
    # BaseTableExtractor Interface Implementation
    # ========================================================================
    
    def detect_table_regions(self, content: Any) -> List[TableRegion]:
        """Detect table regions in PPT/PPTX presentation.
        
        Pass 1: Find all shapes with has_table=True in all slides.
        
        Args:
            content: python-pptx Presentation object
            
        Returns:
            List of TableRegion objects
        """
        prs = self._get_presentation(content)
        if prs is None:
            return []
        
        regions = []
        self._table_info_cache.clear()
        table_idx = 0
        
        try:
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if not hasattr(shape, 'has_table') or not shape.has_table:
                        continue
                    
                    table = shape.table
                    num_rows = len(table.rows)
                    num_cols = len(table.columns)
                    
                    # Skip based on configuration
                    if self._should_skip_table(num_rows, num_cols):
                        continue
                    
                    # Calculate confidence
                    confidence = self._calculate_confidence(table, num_rows, num_cols)
                    
                    # Create region
                    region = TableRegion(
                        start_offset=table_idx,
                        end_offset=table_idx + 1,
                        row_count=num_rows,
                        col_count=num_cols,
                        confidence=confidence,
                    )
                    regions.append(region)
                    
                    # Cache table info for Pass 2
                    shape_id = shape.shape_id if hasattr(shape, 'shape_id') else id(shape)
                    self._table_info_cache[table_idx] = PPTTableRegionInfo(
                        slide_index=slide_idx,
                        shape_id=shape_id,
                        table_ref=table,
                    )
                    
                    table_idx += 1
            
            self._presentation_cache = prs
            self.logger.debug(f"Detected {len(regions)} table regions in PPT")
            
        except Exception as e:
            self.logger.error(f"Error detecting table regions: {e}")
        
        return regions
    
    def extract_table_from_region(
        self, 
        content: Any, 
        region: TableRegion
    ) -> Optional[TableData]:
        """Extract table data from a detected region.
        
        Pass 2: Extract actual table content from region.
        
        Args:
            content: python-pptx Presentation object (or cached)
            region: TableRegion identifying the table
            
        Returns:
            TableData object or None if extraction fails
        """
        table_idx = region.start_offset
        
        # Get cached table info
        table_info = self._table_info_cache.get(table_idx)
        if table_info is None:
            self.logger.warning(f"No cached table info for index {table_idx}")
            return None
        
        table = table_info.table_ref
        if table is None:
            return None
        
        try:
            return self._extract_table_data(table, table_info.slide_index)
        except Exception as e:
            self.logger.error(f"Error extracting table from region: {e}")
            return None
    
    # ========================================================================
    # Helper Methods
    # ========================================================================
    
    def _get_presentation(self, content: Any) -> Any:
        """Get python-pptx Presentation object from various input types."""
        try:
            # If already a Presentation
            if hasattr(content, 'slides'):
                return content
            
            # If bytes, create Presentation from bytes
            if isinstance(content, bytes):
                from io import BytesIO
                from pptx import Presentation
                return Presentation(BytesIO(content))
            
            # If string path
            if isinstance(content, str):
                from pptx import Presentation
                return Presentation(content)
            
            # Use cached if available
            if self._presentation_cache is not None:
                return self._presentation_cache
            
        except Exception as e:
            self.logger.error(f"Failed to get Presentation: {e}")
        
        return None
    
    def _should_skip_table(self, num_rows: int, num_cols: int) -> bool:
        """Check if table should be skipped based on configuration."""
        config = self._config
        
        # 1x1 table
        if config.skip_single_cell_tables and num_rows == 1 and num_cols == 1:
            return True
        
        # Nx1 table (single column)
        if config.skip_single_column_tables and num_cols == 1:
            return True
        
        # 1xN table (single row)
        if config.skip_single_row_tables and num_rows == 1:
            return True
        
        return False
    
    def _calculate_confidence(self, table: Any, num_rows: int, num_cols: int) -> float:
        """Calculate confidence score for table detection."""
        confidence = 0.8  # Base confidence for explicit table
        
        # Bonus for larger tables
        if num_rows >= 3:
            confidence += 0.05
        if num_cols >= 3:
            confidence += 0.05
        
        # Bonus for having content
        has_content = False
        try:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        has_content = True
                        break
                if has_content:
                    break
        except:
            pass
        
        if has_content:
            confidence += 0.05
        
        return min(confidence, 1.0)
    
    def _extract_table_data(self, table: Any, slide_index: int) -> TableData:
        """Extract TableData from python-pptx Table object.
        
        Handles cell merges properly.
        """
        num_rows = len(table.rows)
        num_cols = len(table.columns)
        
        if num_rows == 0 or num_cols == 0:
            return TableData(
                rows=[],
                num_rows=0,
                num_cols=0,
                source_format="pptx",
            )
        
        # Build merge info map
        merge_info = self._build_merge_info(table, num_rows, num_cols)
        
        # Extract cells
        rows_data: List[List[TableCell]] = []
        
        for row_idx in range(num_rows):
            row_cells: List[TableCell] = []
            
            for col_idx in range(num_cols):
                info = merge_info.get((row_idx, col_idx))
                
                if info is None:
                    # Regular cell
                    cell = table.cell(row_idx, col_idx)
                    cell_text = cell.text.strip() if cell.text else ""
                    
                    row_cells.append(TableCell(
                        content=cell_text,
                        row_span=1,
                        col_span=1,
                        is_header=(row_idx == 0 and self._config.treat_first_row_as_header),
                        row_index=row_idx,
                        col_index=col_idx,
                    ))
                
                elif info['skip']:
                    # Merged-away cell - skip or add placeholder
                    row_cells.append(TableCell(
                        content="",
                        row_span=0,
                        col_span=0,
                        is_header=False,
                        row_index=row_idx,
                        col_index=col_idx,
                    ))
                
                else:
                    # Merge origin cell
                    cell = table.cell(row_idx, col_idx)
                    cell_text = cell.text.strip() if cell.text else ""
                    
                    row_cells.append(TableCell(
                        content=cell_text,
                        row_span=info['rowspan'],
                        col_span=info['colspan'],
                        is_header=(row_idx == 0 and self._config.treat_first_row_as_header),
                        row_index=row_idx,
                        col_index=col_idx,
                    ))
            
            rows_data.append(row_cells)
        
        # Calculate column widths
        col_widths = self._calculate_column_widths(table, num_cols)
        
        return TableData(
            rows=rows_data,
            num_rows=num_rows,
            num_cols=num_cols,
            has_header=self._config.treat_first_row_as_header,
            start_offset=slide_index,
            end_offset=slide_index,
            source_format="pptx",
            metadata={"slide_index": slide_index},
            col_widths_percent=col_widths,
        )
    
    def _build_merge_info(
        self, 
        table: Any, 
        num_rows: int, 
        num_cols: int
    ) -> Dict[Tuple[int, int], Dict[str, Any]]:
        """Build merge information map for table.
        
        Returns dict mapping (row_idx, col_idx) -> {'rowspan', 'colspan', 'skip'}
        """
        merge_info: Dict[Tuple[int, int], Dict[str, Any]] = {}
        
        for row_idx in range(num_rows):
            for col_idx in range(num_cols):
                if (row_idx, col_idx) in merge_info:
                    continue
                
                try:
                    cell = table.cell(row_idx, col_idx)
                    cell_merge = self._get_cell_merge_info(
                        cell, table, row_idx, col_idx, num_rows, num_cols
                    )
                    
                    rowspan = cell_merge['rowspan']
                    colspan = cell_merge['colspan']
                    
                    if rowspan == 0 and colspan == 0:
                        # This cell is spanned by another cell
                        merge_info[(row_idx, col_idx)] = {'skip': True}
                        continue
                    
                    if rowspan > 1 or colspan > 1:
                        # Merge origin
                        merge_info[(row_idx, col_idx)] = {
                            'rowspan': rowspan,
                            'colspan': colspan,
                            'skip': False,
                        }
                        
                        # Mark merged-away cells
                        for r in range(row_idx, min(row_idx + rowspan, num_rows)):
                            for c in range(col_idx, min(col_idx + colspan, num_cols)):
                                if r != row_idx or c != col_idx:
                                    merge_info[(r, c)] = {'skip': True}
                
                except Exception as e:
                    self.logger.debug(f"Error getting merge info at ({row_idx}, {col_idx}): {e}")
        
        return merge_info
    
    def _get_cell_merge_info(
        self,
        cell: Any,
        table: Any,
        row_idx: int,
        col_idx: int,
        num_rows: int,
        num_cols: int,
    ) -> Dict[str, int]:
        """Get merge info for a single cell.
        
        Tries multiple methods:
        1. python-pptx built-in attributes (is_merge_origin, span_width, span_height)
        2. XML parsing (gridSpan, rowSpan)
        3. Reference comparison fallback
        """
        rowspan = 1
        colspan = 1
        
        try:
            # Method 1: python-pptx built-in attributes
            if hasattr(cell, 'is_merge_origin') and cell.is_merge_origin:
                if hasattr(cell, 'span_height'):
                    rowspan = cell.span_height
                if hasattr(cell, 'span_width'):
                    colspan = cell.span_width
                return {'rowspan': rowspan, 'colspan': colspan}
            
            # Check if this cell is spanned by another
            if hasattr(cell, 'is_spanned') and cell.is_spanned:
                return {'rowspan': 0, 'colspan': 0}
            
            # Method 2: XML parsing
            tc = cell._tc
            
            grid_span = tc.get('gridSpan')
            if grid_span:
                colspan = int(grid_span)
            
            row_span_attr = tc.get('rowSpan')
            if row_span_attr:
                rowspan = int(row_span_attr)
            
            # Method 3: Reference comparison fallback
            if colspan == 1:
                colspan = self._detect_colspan_by_reference(table, row_idx, col_idx, num_cols)
            
            if rowspan == 1:
                rowspan = self._detect_rowspan_by_reference(table, row_idx, col_idx, num_rows)
        
        except Exception as e:
            self.logger.debug(f"Error getting cell merge info: {e}")
        
        return {'rowspan': rowspan, 'colspan': colspan}
    
    def _detect_colspan_by_reference(
        self, 
        table: Any, 
        row_idx: int, 
        col_idx: int, 
        num_cols: int
    ) -> int:
        """Detect colspan by comparing cell references."""
        colspan = 1
        try:
            current_cell = table.cell(row_idx, col_idx)
            
            for c in range(col_idx + 1, num_cols):
                next_cell = table.cell(row_idx, c)
                if next_cell._tc is current_cell._tc:
                    colspan += 1
                else:
                    break
        except:
            pass
        
        return colspan
    
    def _detect_rowspan_by_reference(
        self, 
        table: Any, 
        row_idx: int, 
        col_idx: int, 
        num_rows: int
    ) -> int:
        """Detect rowspan by comparing cell references."""
        rowspan = 1
        try:
            current_cell = table.cell(row_idx, col_idx)
            
            for r in range(row_idx + 1, num_rows):
                next_cell = table.cell(r, col_idx)
                if next_cell._tc is current_cell._tc:
                    rowspan += 1
                else:
                    break
        except:
            pass
        
        return rowspan
    
    def _calculate_column_widths(self, table: Any, num_cols: int) -> List[float]:
        """Calculate column widths as percentages."""
        if num_cols == 0:
            return []
        
        try:
            # Try to get actual column widths
            total_width = 0
            widths = []
            
            for col in table.columns:
                width = col.width.emu if hasattr(col.width, 'emu') else col.width
                widths.append(width)
                total_width += width
            
            if total_width > 0:
                return [(w / total_width) * 100 for w in widths]
        except:
            pass
        
        # Default: equal widths
        return [100.0 / num_cols] * num_cols
    
    # ========================================================================
    # Convenience Methods
    # ========================================================================
    
    def is_simple_table(self, table: Any) -> bool:
        """Check if table is a simple table (1xN or Nx1).
        
        Simple tables are often used as text containers, not real tables.
        """
        try:
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            
            if num_rows == 1 or num_cols == 1:
                return True
            
            return False
        except:
            return False
    
    def extract_tables_from_slide(
        self, 
        slide: Any, 
        slide_index: int = 0
    ) -> List[TableData]:
        """Extract all tables from a single slide.
        
        Convenience method for slide-by-slide processing.
        """
        tables = []
        
        try:
            for shape in slide.shapes:
                if not hasattr(shape, 'has_table') or not shape.has_table:
                    continue
                
                table = shape.table
                num_rows = len(table.rows)
                num_cols = len(table.columns)
                
                if self._should_skip_table(num_rows, num_cols):
                    continue
                
                table_data = self._extract_table_data(table, slide_index)
                if table_data and table_data.is_valid(self.config.min_rows, self.config.min_cols):
                    tables.append(table_data)
        
        except Exception as e:
            self.logger.error(f"Error extracting tables from slide: {e}")
        
        return tables


# ============================================================================
# Convenience Functions
# ============================================================================

def extract_tables_from_presentation(
    presentation: Any,
    config: Optional[PPTTableExtractorConfig] = None,
) -> List[TableData]:
    """Extract all tables from a presentation.
    
    Convenience function for simple usage.
    
    Args:
        presentation: python-pptx Presentation object
        config: Optional extraction configuration
        
    Returns:
        List of TableData objects
    """
    extractor = PPTTableExtractor(config)
    return extractor.extract_tables(presentation)


def extract_table_from_shape(
    shape: Any,
    slide_index: int = 0,
    config: Optional[PPTTableExtractorConfig] = None,
) -> Optional[TableData]:
    """Extract table data from a single shape.
    
    Args:
        shape: python-pptx Shape object with has_table=True
        slide_index: Index of the slide containing this shape
        config: Optional extraction configuration
        
    Returns:
        TableData or None
    """
    if not hasattr(shape, 'has_table') or not shape.has_table:
        return None
    
    extractor = PPTTableExtractor(config)
    return extractor._extract_table_data(shape.table, slide_index)


# ============================================================================
# Backward Compatibility
# ============================================================================

def is_simple_table(table: Any) -> bool:
    """Check if table is simple (1xN or Nx1).
    
    Backward compatibility wrapper.
    """
    extractor = PPTTableExtractor()
    return extractor.is_simple_table(table)


def debug_table_structure(table: Any) -> None:
    """Debug table structure for merge info analysis.
    
    Backward compatibility wrapper from ppt_table.py.
    
    Args:
        table: python-pptx Table object
    """
    logger.debug("=== Table Structure Debug ===")
    logger.debug(f"Rows: {len(table.rows)}, Cols: {len(table.columns)}")
    
    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            try:
                cell = table.cell(row_idx, col_idx)
                tc = cell._tc
                
                # XML attributes
                grid_span = tc.get('gridSpan', '1')
                row_span = tc.get('rowSpan', '1')
                
                # python-pptx attributes
                is_merge_origin = getattr(cell, 'is_merge_origin', None)
                is_spanned = getattr(cell, 'is_spanned', None)
                span_width = getattr(cell, 'span_width', None)
                span_height = getattr(cell, 'span_height', None)
                
                text_preview = cell.text[:20] if cell.text else ""
                
                logger.debug(
                    f"[{row_idx},{col_idx}] "
                    f"text='{text_preview}' "
                    f"gridSpan={grid_span} rowSpan={row_span} "
                    f"is_merge_origin={is_merge_origin} "
                    f"is_spanned={is_spanned} "
                    f"span_width={span_width} span_height={span_height}"
                )
            except Exception as e:
                logger.debug(f"[{row_idx},{col_idx}] Error: {e}")
    
    logger.debug("=== End Debug ===")


# ============================================================================
# Exports
# ============================================================================

__all__ = [
    # Main classes
    'PPTTableExtractor',
    'PPTTableExtractorConfig',
    'PPTTableRegionInfo',
    # Convenience functions
    'extract_tables_from_presentation',
    'extract_table_from_shape',
    # Backward compatibility
    'is_simple_table',
    'debug_table_structure',
]
